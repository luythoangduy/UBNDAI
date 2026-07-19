"""Sinh deck pitching UBNDAI theo khung trong pitching.pdf (Sequoia 12 slide).

Nguyên tắc áp dụng từ tài liệu:
- Outcome first, technology second, proof always.
- Không mở đầu bằng công nghệ; mở đầu bằng nỗi đau khách hàng.
- Demo là câu chuyện có nhân vật, không phải walkthrough nút bấm.
- "Nếu dùng số liệu cụ thể, phải luôn có nguồn xác thực" → ô nào chưa có số
  liệu thật thì đánh dấu đỏ để đội tự điền, KHÔNG bịa.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0D, 0x31, 0x58)
BLUE = RGBColor(0x17, 0x6F, 0xA6)
TEAL = RGBColor(0x16, 0x7C, 0x78)
CREAM = RGBColor(0xF3, 0xF1, 0xEA)
INK = RGBColor(0x1B, 0x2A, 0x3D)
MUTED = RGBColor(0x5B, 0x70, 0x85)
ALERT = RGBColor(0xC2, 0x5A, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h):
    return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def text(tb, runs, align=PP_ALIGN.LEFT, space_after=6):
    """runs: list of (text, size, bold, color) — mỗi phần tử là một đoạn."""
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for content, size, bold, color in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = content
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
    return tb


def bar(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def card(s, x, y, w, h, title, body, accent=BLUE, title_size=15, body_size=12):
    bar(s, x, y, w, h, WHITE)
    bar(s, x, y, 0.055, h, accent)
    tb = box(s, x + 0.25, y + 0.16, w - 0.45, h - 0.3)
    text(tb, [(title, title_size, True, NAVY), (body, body_size, False, MUTED)], space_after=4)


def header(s, eyebrow, title, sub=None):
    tb = box(s, 0.75, 0.45, 11.9, 1.5)
    runs = [(eyebrow, 11, True, BLUE), (title, 33, True, NAVY)]
    if sub:
        runs.append((sub, 14, False, MUTED))
    text(tb, runs, space_after=3)


def note(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def needs_data(s, x, y, w, h, what):
    """Ô đánh dấu rõ: chưa có số liệu thật, đội phải tự điền."""
    bar(s, x, y, w, h, RGBColor(0xFD, 0xF1, 0xE7))
    bar(s, x, y, 0.055, h, ALERT)
    tb = box(s, x + 0.25, y + 0.14, w - 0.45, h - 0.28)
    text(tb, [("⚠ CẦN SỐ LIỆU THẬT", 10, True, ALERT), (what, 12, False, INK)], space_after=4)


# ── 1. Title ─────────────────────────────────────────────────────────────
s = slide(NAVY)
bar(s, 0, 6.9, 13.333, 0.6, BLUE)
tb = box(s, 1.0, 2.0, 11.3, 3.2)
text(tb, [
    ("CỔNG DỊCH VỤ CÔNG · TRỢ LÝ AI", 13, True, RGBColor(0x8F, 0xC7, 0xE8)),
    ("UBNDAI", 60, True, WHITE),
    ("Hồ sơ hành chính đúng ngay từ lần nộp đầu tiên.", 22, False, RGBColor(0xD5, 0xE6, 0xF2)),
], space_after=10)
tb = box(s, 1.0, 5.45, 11.3, 1.0)
text(tb, [("Chúng tôi giúp NGƯỜI DÂN nộp đúng hồ sơ ngay lần đầu bằng TRỢ LÝ AI TIỀN KIỂM "
           "DỰA TRÊN NGUỒN PHÁP LÝ CHÍNH THỨC, thay vì ĐI LẠI NHIỀU LẦN VÌ THIẾU GIẤY TỜ.",
           14, False, RGBColor(0xA9, 0xCB, 0xE3))])
note(s, "0:00–0:30 | HOOK (10%). Đọc chậm one-liner. Không nhắc chữ 'LLM', 'RAG' ở slide này.")

# ── 2. Hook ──────────────────────────────────────────────────────────────
s = slide()
header(s, "MỞ ĐẦU", "Người dân không sợ thủ tục.\nHọ sợ phải quay lại lần thứ ba.")
for i, (t, b) in enumerate([
    ("Nộp hồ sơ", "Xếp hàng, nộp, chờ."),
    ("Bị trả lại", "Thiếu một tờ giấy. Sai một ô. Không ai nói trước."),
    ("Quay lại từ đầu", "Nghỉ làm thêm một buổi. Lặp lại vòng lặp."),
]):
    card(s, 0.75 + i * 4.05, 2.6, 3.8, 1.55, t, b, accent=ALERT if i == 1 else BLUE)
needs_data(s, 0.75, 4.5, 11.8, 1.15,
           "Một con số mở đầu, kèm nguồn: tỷ lệ hồ sơ bị trả lại / số lượt đi lại trung bình "
           "cho một thủ tục ở xã. Không có số thật thì kể một câu chuyện có thật thay thế.")
tb = box(s, 0.75, 5.85, 11.8, 0.7)
text(tb, [("Nỗi đau không nằm ở việc điền form. Nó nằm ở chỗ không ai kiểm tra giúp bạn trước khi quá muộn.",
           15, True, NAVY)])
note(s, "0:30–1:15 | PROBLEM (15%). Mở bằng nỗi đau, KHÔNG mở bằng công nghệ (lỗi #1 trong tài liệu).")

# ── 3. Customer ──────────────────────────────────────────────────────────
s = slide()
header(s, "KHÁCH HÀNG", "Ba nhóm, một vòng lặp chung")
for i, (t, b, c) in enumerate([
    ("Người dân", "Làm khai sinh, kết hôn, tạm trú, căn cước, giấy phép xây dựng. "
                  "Không rành thuật ngữ pháp lý, không biết mình thiếu gì.", BLUE),
    ("Cán bộ tiếp nhận (UBND cấp xã)", "Mỗi hồ sơ sai là một lần giải thích lại, một lần nhập lại. "
                                       "Thời gian dồn vào việc lẽ ra máy làm được.", TEAL),
    ("Cơ quan quản lý", "Cần biết hồ sơ hay sai ở đâu, thủ tục nào nghẽn — hiện gần như không có dữ liệu.", NAVY),
]):
    card(s, 0.75, 2.5 + i * 1.45, 11.8, 1.28, t, b, accent=c, body_size=12)
needs_data(s, 0.75, 6.05, 11.8, 0.85,
           "Đã phỏng vấn/khảo sát bao nhiêu người dân và cán bộ? Ở xã/phường nào? "
           "Đây là câu BA Judge hỏi thẳng ở checkpoint 36h.")
note(s, "Người trả tiền (cơ quan) khác người dùng chính (người dân) — phải nói rõ ở slide Business Model.")

# ── 4. Solution ──────────────────────────────────────────────────────────
s = slide()
header(s, "GIẢI PHÁP", "Chuyển việc kiểm tra hồ sơ\nra TRƯỚC khi nộp")
tb = box(s, 0.75, 2.35, 11.8, 0.6)
text(tb, [("Không phải chatbot trả lời thủ tục. Là lớp tiền kiểm đứng giữa người dân và quầy tiếp nhận.",
           15, True, NAVY)])
for i, (t, b) in enumerate([
    ("1 · Hỏi bằng lời thường", "Người dân mô tả việc mình cần. Hệ thống nhận diện đúng thủ tục "
                                "và hỏi lại đúng vài câu để cá nhân hoá checklist."),
    ("2 · Tự đọc giấy tờ", "Chụp ảnh giấy tờ, AI bóc dữ liệu và điền vào tờ khai. "
                           "Trường nào độ tin cậy thấp thì đánh dấu để người kiểm."),
    ("3 · Kiểm trước khi nộp", "Bộ quy tắc nghiệp vụ soát hồ sơ và chỉ ra chỗ sai — "
                               "trước khi người dân rời khỏi nhà."),
]):
    card(s, 0.75 + i * 4.05, 3.15, 3.8, 2.15, t, b, body_size=12)
tb = box(s, 0.75, 5.6, 11.8, 1.2)
text(tb, [("Kết quả: bớt một vòng đi lại. Cán bộ nhận hồ sơ đã sạch. Cơ quan có dữ liệu về chỗ hay sai.",
           15, True, TEAL)])
note(s, "1:15–2:00 | SOLUTION (15%). Nói KẾT QUẢ trước, công nghệ để dành slide sau.")

# ── 5. Demo story ────────────────────────────────────────────────────────
s = slide()
header(s, "DEMO", "Chị Lan, 34 tuổi, cần làm khai sinh cho con",
       "Demo là câu chuyện — không phải đi qua từng nút bấm.")
for i, (t, b, c) in enumerate([
    ("1. Bối cảnh", "Chị nhắn: “tôi muốn đăng ký khai sinh cho con”. Không biết cần giấy gì, "
                    "cũng không biết nộp ở đâu.", MUTED),
    ("2. Can thiệp", "Trợ lý nhận đúng thủ tục, hỏi 2 câu (bé sinh ở cơ sở y tế? cha mẹ đã kết hôn?), "
                     "rồi dựng checklist riêng cho trường hợp của chị.", BLUE),
    ("3. Khoảnh khắc Aha", "Chị chụp giấy chứng sinh. Tờ khai tự điền. Hệ thống báo ngay: "
                           "“còn thiếu mục X” — trước khi chị ra khỏi nhà.", ALERT),
    ("4. Lợi thế", "Mỗi câu trả lời đều kèm nguồn: mã thủ tục, cơ quan, biểu mẫu tải thẳng "
                   "từ Cổng Dịch vụ công.", TEAL),
]):
    card(s, 0.75 + i * 3.02, 2.75, 2.85, 2.5, t, b, accent=c, title_size=14, body_size=11)
tb = box(s, 0.75, 5.5, 11.8, 1.1)
text(tb, [("Chốt demo bằng một câu: “Chị Lan đi một lần, không đi ba lần.”", 16, True, NAVY)])
note(s, "2:00–3:00 | DEMO (20%) — phần nặng điểm nhất. Dừng lâu ở bước 3. "
        "Đừng giải thích kiến trúc trong lúc demo.")

# ── 6. AI-Native ─────────────────────────────────────────────────────────
s = slide()
header(s, "VÌ SAO CẦN AI", "AI là động cơ phía sau kết quả — không phải khẩu hiệu")
for i, (t, b) in enumerate([
    ("Hiểu ý định nói thường", "Người dân không tra cứu bằng mã thủ tục. Họ nói “làm giấy cho con mới đẻ”. "
                               "Hệ thống phải tự khớp về đúng thủ tục trong danh mục."),
    ("Đọc giấy tờ ảnh chụp", "Ảnh chụp bằng điện thoại, nghiêng, mờ. AI bóc trường dữ liệu và "
                             "tự đánh dấu chỗ nó không chắc."),
    ("Checklist theo trường hợp", "Cùng thủ tục giấy phép xây dựng: danh mục có 5 thành phần, "
                                  "nhưng tuỳ câu trả lời chỉ còn 3 — 4 — 5 mục thật sự phải chuẩn bị."),
]):
    card(s, 0.75, 2.5 + i * 1.42, 11.8, 1.25, t, b, body_size=12)
bar(s, 0.75, 6.0, 11.8, 0.95, RGBColor(0xE9, 0xF3, 0xF0))
tb = box(s, 1.0, 6.15, 11.3, 0.7)
text(tb, [("Ranh giới cố ý: AI chỉ được cảnh báo. Kết luận “hồ sơ sai” chỉ do bộ quy tắc nghiệp vụ đưa ra. "
           "Quyết định cuối cùng vẫn thuộc về cán bộ.", 13, True, TEAL)])
note(s, "Đây là slide chống lại lỗi 'AI của em rất thông minh'. Luôn quy về kết quả và ranh giới an toàn.")

# ── 7. Proof ─────────────────────────────────────────────────────────────
s = slide()
header(s, "BẰNG CHỨNG", "Những gì đã chạy được, kiểm chứng được")
for i, (t, b) in enumerate([
    ("Sản phẩm chạy thật", "Luồng đầy đủ người dân → cán bộ: hỏi đáp, sinh checklist, "
                           "OCR điền tờ khai, kiểm tra, nộp, thẩm tra, ra quyết định."),
    ("Nguồn truy vết được", "5 thủ tục · 17 luật kiểm tra khai báo · 75 endpoint. Kiểm thử tự động "
                            "bắt buộc mỗi trích dẫn mở đúng mã thủ tục quốc gia của chính nó."),
    ("Kỷ luật kỹ thuật", "335 bài kiểm thử, CI xanh. Cổng CI chặn merge khi hệ thống chốt nhầm "
                         "thủ tục — ngưỡng cho phép: 0. Đo được, không phải lời hứa."),
]):
    card(s, 0.75 + i * 4.05, 2.5, 3.8, 2.35, t, b, body_size=11.5)
bar(s, 0.75, 5.15, 11.8, 1.0, WHITE)
tb = box(s, 1.0, 5.28, 11.3, 0.8)
text(tb, [("Chốt nhầm thủ tục: 0/60 trên bộ đánh giá giữ riêng. Khi không chắc, hệ thống hỏi lại "
           "chứ không đoán — vì chốt nhầm thủ tục làm sai toàn bộ checklist phía sau.", 13, True, NAVY)])
needs_data(s, 0.75, 6.3, 11.8, 0.75,
           "Traction thật: số người dùng thử, số hồ sơ chạy qua hệ thống, phản hồi của cán bộ.")
note(s, "3:00–4:00 | BIZ & TRACTION (20%). Nếu chưa có traction thị trường, đừng giả vờ — "
        "đưa bằng chứng kỹ thuật + cam kết pilot.")

# ── 7b. An toàn đo được ──────────────────────────────────────────────────
s = slide()
header(s, "AN TOÀN", "Không hứa an toàn. Đo nó.")
for i, (t, b) in enumerate([
    ("Chốt nhầm thủ tục: 0/60",
     "Trên bộ đánh giá giữ riêng — viết TRƯỚC khi tinh chỉnh và không dùng để tinh chỉnh. "
     "10 ca không nhận ra đều lùi về hỏi lại, không ca nào đoán bừa."),
    ("AI không thể kết luận “hồ sơ sai”",
     "Chặn ở tầng kiểu dữ liệu, không phải lời dặn trong prompt. Thử tạo cũng ném lỗi. "
     "Chỉ bộ quy tắc khai báo — thứ cán bộ đọc và ký duyệt được — mới có thẩm quyền đó."),
    ("Bộ đánh giá giữ riêng đã bắt được lỗi thật",
     "Ba lỗi khiến hệ thống chốt nhầm thủ tục, trong đó hai lỗi đặc thù tiếng Việt: bỏ dấu làm "
     "“chứng minh” trùng “chung mình”, và “cưới” trùng “cuối”. Đã sửa, có kiểm thử khoá lại."),
]):
    card(s, 0.75, 2.5 + i * 1.42, 11.8, 1.25, t, b, body_size=11.5, accent=TEAL)
bar(s, 0.75, 6.0, 11.8, 0.95, RGBColor(0xE9, 0xF3, 0xF0))
tb = box(s, 1.0, 6.15, 11.3, 0.7)
text(tb, [("Chốt nhầm thủ tục làm sai TOÀN BỘ checklist phía sau — người dân chuẩn bị đúng một bộ "
           "giấy tờ cho sai một thủ tục. Vì vậy nó là cổng chặn merge, không phải chỉ số cải thiện dần.",
           12.5, True, TEAL)])
note(s, "Slide dành cho tiêu chí AI Safety & Trust. Nếu bị hỏi 'độ chính xác bao nhiêu', "
        "trả lời bằng chốt nhầm 0/60 trước, rồi mới nói độ phủ 60% trên bộ giữ riêng — "
        "tự nêu điểm yếu trước khi giám khảo nêu.")

# ── 8. Moat ──────────────────────────────────────────────────────────────
s = slide()
header(s, "LỢI THẾ", "Vì sao khó sao chép")
for i, (t, b) in enumerate([
    ("Dữ liệu thủ tục có kỷ luật nguồn", "Không phải bộ prompt. Là danh mục thủ tục chuẩn hoá, "
                                         "mỗi mục truy vết được về văn bản công bố, có quy trình đồng bộ khi nguồn đổi."),
    ("Tách bạch AI và luật", "Luật kiểm tra viết dưới dạng khai báo, tách khỏi mô hình. "
                             "Thêm thủ tục mới không cần huấn luyện lại, và không làm AI 'chế' ra quy định."),
    ("Vòng phản hồi từ cán bộ", "Mỗi lần cán bộ sửa kết quả AI là một nhãn dữ liệu. "
                                "Càng dùng, checklist càng khớp thực tế địa phương."),
]):
    card(s, 0.75, 2.5 + i * 1.5, 11.8, 1.32, t, b, body_size=12)
tb = box(s, 0.75, 6.2, 11.8, 0.7)
text(tb, [("Đối thủ dựng được con chatbot trong một tuần. Dựng được lớp truy vết nguồn thì không.",
           15, True, NAVY)])
note(s, "Trả lời câu 'MOAT' trong 7 lớp lọc của nhà đầu tư.")

# ── 9. Market ────────────────────────────────────────────────────────────
s = slide()
header(s, "THỊ TRƯỜNG", "Quy mô cơ hội")
needs_data(s, 0.75, 2.4, 11.8, 1.5,
           "TAM / SAM / SOM. Gợi ý nguồn công khai để lấy số: số đơn vị hành chính cấp xã sau sáp nhập; "
           "tổng số hồ sơ TTHC tiếp nhận/năm trên Cổng Dịch vụ công Quốc gia; ngân sách chuyển đổi số cấp tỉnh.")
needs_data(s, 0.75, 4.05, 5.8, 1.5, "Vì sao là LÚC NÀY: chính sách chuyển đổi số, "
                                    "phân cấp về cấp xã, dữ liệu dân cư đã liên thông.")
needs_data(s, 6.75, 4.05, 5.8, 1.5, "Đã có xã/phường nào đồng ý thử? Tên đơn vị + người liên hệ "
                                    "có sức nặng hơn mọi con số ước lượng.")
tb = box(s, 0.75, 5.8, 11.8, 0.9)
text(tb, [("Tài liệu pitching nhấn: “Nếu dùng số liệu cụ thể, phải luôn có nguồn xác thực.” "
           "Ba ô trên để trống có chủ đích — điền bằng số thật, đừng điền bằng số đoán.", 13, True, ALERT)])
note(s, "Đây là điểm yếu lớn nhất hiện tại của đội (tiêu chí Business Viability = 20 điểm).")

# ── 10. Business model ───────────────────────────────────────────────────
s = slide()
header(s, "MÔ HÌNH KINH DOANH", "Ai trả tiền, trả cho cái gì")
for i, (t, b) in enumerate([
    ("Người dùng cuối", "Người dân — miễn phí. Đây là dịch vụ công, thu phí người dân là phản tác dụng."),
    ("Người trả tiền", "Cơ quan nhà nước (UBND cấp xã/huyện/tỉnh) — trả cho việc giảm tải quầy tiếp nhận "
                       "và có số liệu điều hành."),
    ("Hình thức", "Thuê bao theo đơn vị hành chính + phí triển khai ban đầu (chuẩn hoá danh mục thủ tục địa phương)."),
]):
    card(s, 0.75, 2.5 + i * 1.35, 11.8, 1.18, t, b, body_size=12)
tbc = box(s, 0.75, 6.62, 11.8, 0.58)
text(tbc, [("Chi phí AI mỗi hồ sơ: ~410 đ (0,016 USD) — tính từ giá token thật và tỉ lệ 2,17 ký tự/token "
            "đo trên chính prompt của hệ thống. Ở mức đó, miễn phí cho người dân là bền vững.",
            12, True, TEAL)])
needs_data(s, 0.75, 7.22, 11.8, 0.42, "Còn thiếu: đơn giá thuê bao và điểm hoà vốn.")
tb = box(s, 0.75, 6.0, 11.8, 0.6)
text(tb, [("Lưu ý: mô hình mua sắm công có quy trình riêng — cần nêu rõ đường vào (thí điểm → nghiệm thu → mở rộng).",
           12.5, False, MUTED)])
note(s, "Trả lời lỗi #4: 'Không rõ ai trả tiền'.")

# ── 11. Team ─────────────────────────────────────────────────────────────
s = slide()
header(s, "ĐỘI NGŨ", "Ai thực thi")
needs_data(s, 0.75, 2.5, 11.8, 2.0,
           "Tên, vai trò, và VÌ SAO chính người này làm được việc này (kinh nghiệm liên quan: "
           "AI/RAG, backend, khu vực công, thiết kế). Nhà đầu tư đọc slide Team để đánh giá năng lực thực thi.")
for i, (t, b) in enumerate([
    ("Đã chứng minh được", "Dựng xong luồng đầy đủ hai phía trong thời gian hackathon, "
                           "có kiểm thử tự động và kỷ luật nguồn dữ liệu."),
    ("Đang cần bổ sung", "Người có quan hệ triển khai khu vực công và kinh nghiệm mua sắm công."),
]):
    card(s, 0.75 + i * 6.05, 4.7, 5.8, 1.5, t, b, accent=TEAL if i == 0 else ALERT, body_size=12)
note(s, "Thành thật về chỗ còn thiếu tạo niềm tin hơn là tỏ ra đủ mọi thứ.")

# ── 12. Roadmap ──────────────────────────────────────────────────────────
s = slide()
header(s, "LỘ TRÌNH", "Từ demo đến triển khai thật")
stages = [
    ("Giai đoạn 1\nThí điểm", "1 xã/phường · 5 thủ tục hiện có\nĐo: tỷ lệ hồ sơ đạt ngay lần đầu", BLUE),
    ("Giai đoạn 2\nMở rộng thủ tục", "Bổ sung nhóm thủ tục hộ tịch – cư trú – đất đai\nChuẩn hoá quy trình đồng bộ nguồn", TEAL),
    ("Giai đoạn 3\nNhân rộng", "Nhiều đơn vị trong cùng tỉnh\nBảng điều hành cho cơ quan quản lý", NAVY),
]
for i, (t, b, c) in enumerate(stages):
    card(s, 0.75 + i * 4.05, 2.6, 3.8, 2.1, t, b, accent=c, body_size=11.5)
needs_data(s, 0.75, 5.0, 11.8, 1.0,
           "Mốc thời gian cụ thể cho từng giai đoạn + chỉ số nghiệm thu đã thống nhất với đơn vị thí điểm.")
tb = box(s, 0.75, 6.2, 11.8, 0.7)
text(tb, [("Chỉ số Bắc Đẩu: tỷ lệ hồ sơ được tiếp nhận ngay lần nộp đầu tiên.", 15, True, NAVY)])
note(s, "4:00–5:00 | VISION & ASK (20%).")

# ── 13. Ask ──────────────────────────────────────────────────────────────
s = slide(NAVY)
tb = box(s, 0.9, 1.1, 11.5, 1.2)
text(tb, [("LỜI ĐỀ NGHỊ", 13, True, RGBColor(0x8F, 0xC7, 0xE8)),
          ("Chúng tôi cần gì để đi tiếp", 38, True, WHITE)], space_after=6)
for i, (t, b) in enumerate([
    ("Một đơn vị thí điểm", "Giới thiệu tới một UBND cấp xã sẵn sàng chạy thử 5 thủ tục trong 3 tháng."),
    ("Cố vấn khu vực công", "Người hiểu quy trình công bố thủ tục và mua sắm công."),
    ("Nguồn lực vận hành", "Đủ để duy trì hạ tầng và đồng bộ danh mục thủ tục trong thời gian thí điểm."),
]):
    bar(s, 0.9 + i * 3.95, 2.75, 3.7, 1.9, RGBColor(0x15, 0x3E, 0x6B))
    bar(s, 0.9 + i * 3.95, 2.75, 3.7, 0.06, BLUE)
    tb = box(s, 1.1 + i * 3.95, 2.95, 3.3, 1.6)
    text(tb, [(t, 16, True, WHITE), (b, 12, False, RGBColor(0xB9, 0xD3, 0xE8))], space_after=5)
needs_data(s, 0.9, 5.0, 11.5, 0.95,
           "Điền theo công thức trong tài liệu: “Chúng tôi cần [nguồn lực] để đạt [mục tiêu] trong [thời gian].” "
           "Phải là một câu cụ thể, có con số.")
tb = box(s, 0.9, 6.2, 11.5, 0.8)
text(tb, [("Mục tiêu 3 tháng: chứng minh tỷ lệ hồ sơ đạt ngay lần đầu tăng rõ rệt tại một đơn vị thật.",
           15, True, RGBColor(0xA9, 0xCB, 0xE3))])
note(s, "Lỗi #5 trong tài liệu là 'thiếu Ask rõ ràng'. Câu Ask phải cụ thể, có số, có thời hạn.")

# ── 14. Closing ──────────────────────────────────────────────────────────
s = slide()
tb = box(s, 1.4, 2.5, 10.5, 2.4)
text(tb, [("Người dân không cần hiểu AI.", 34, True, MUTED),
          ("Họ cần đi một lần là xong.", 44, True, NAVY)],
     align=PP_ALIGN.CENTER, space_after=14)
bar(s, 5.4, 5.15, 2.5, 0.05, BLUE)
tb = box(s, 1.4, 5.45, 10.5, 1.0)
text(tb, [("UBNDAI · Trợ lý thủ tục hành chính", 16, True, NAVY),
          ("ubndai.vercel.app", 13, False, MUTED)],
     align=PP_ALIGN.CENTER, space_after=4)
note(s, "Kết bằng lợi ích người dùng, không kết bằng công nghệ.")

out = r"C:\AI_Chalenge\UBNDAI\docs\UBNDAI-pitch-deck.pptx"
prs.save(out)
print("Đã tạo:", out)
print("Số slide:", len(prs.slides.__iter__.__self__._sldIdLst))
